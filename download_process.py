import cv2
import os
from pytube import YouTube
from PIL import Image
import imagehash
from pptx import Presentation
from pptx.util import Inches




def download_and_process_video(video_url, output_folder, ppt_name, resolution="720p"):
    # Download YouTube video
    video_path = download_youtube_video(video_url, output_folder, resolution)

    # Process video frames and create PowerPoint presentation
    capture_frames_with_transition(video_path, output_folder, ppt_name)






def download_youtube_video(video_url, output_folder, resolution="1080p"):
    yt = YouTube(video_url)
    video_stream = yt.streams.filter(file_extension="mp4", res=resolution).first()
    video_path = video_stream.download(output_folder)
    return video_path

def dhash_difference(hash1, hash2):
    return bin(int(hash1, 16) ^ int(hash2, 16)).count('1')

def calculate_frame_difference(prev_frame, current_frame):
    # Convert frames to PIL Image
    prev_image = Image.fromarray(cv2.cvtColor(prev_frame, cv2.COLOR_BGR2RGB))
    current_image = Image.fromarray(cv2.cvtColor(current_frame, cv2.COLOR_BGR2RGB))

    # Compute dHash for the images
    prev_hash = str(imagehash.dhash(prev_image))
    current_hash = str(imagehash.dhash(current_image))

    # Calculate the Hamming distance between the hashes
    diff = dhash_difference(prev_hash, current_hash)

    return diff

def capture_frames_with_transition(video_path, output_folder, ppt_name):
    # Open the video file
    cap = cv2.VideoCapture(video_path)

    # Check if the video file opened successfully
    if not cap.isOpened():
        print("Error opening video file")
        return

    # Create the output folder if it doesn't exist
    os.makedirs(output_folder, exist_ok=True)

    # Get the frames per second (fps) of the video
    fps = cap.get(cv2.CAP_PROP_FPS)

    # Calculate the frame interval in terms of the number of frames (constant interval of 1 second)
    frame_interval = int(fps * 1)

    # Initialize variables for the previous frame and its timestamp
    prev_frame = None
    prev_timestamp = 0

    # Initialize a PowerPoint presentation
    presentation = Presentation()

    # Flag to check if the first frame has been processed
    first_frame_processed = False

    while True:
        # Read a frame from the video
        ret, frame = cap.read()

        # Break the loop if the video has ended
        if not ret:
            break

        # Get the current timestamp
        current_timestamp = cap.get(cv2.CAP_PROP_POS_MSEC) / 1000.0

        # For the first frame (0 seconds), add it to the PowerPoint presentation without difference check
        if not first_frame_processed:
            output_filename = f"{output_folder}/frame_0000.png"
            cv2.imwrite(output_filename, frame)
            print(f"Saved {output_filename}")

            # Add the frame to the PowerPoint presentation
            slide = presentation.slides.add_slide(presentation.slide_layouts[5])

            # Calculate the aspect ratio of the frame
            img_width, img_height = Image.open(output_filename).size
            aspect_ratio = img_width / img_height

            # Set the maximum width to almost fill the slide width
            max_picture_width = Inches(9.5)  # Adjust as needed based on your preference

            # Set the maximum height based on the aspect ratio
            max_picture_height = max_picture_width / aspect_ratio

            # Calculate the centered position
            slide_width, slide_height = presentation.slide_width, presentation.slide_height
            x_position = (slide_width - max_picture_width) / 2
            y_position = (slide_height - max_picture_height) / 2

            # Add the picture to the slide
            slide.shapes.add_picture(output_filename, x_position, y_position, width=max_picture_width, height=max_picture_height)

            # Set the flag to indicate that the first frame has been processed
            first_frame_processed = True

        # If it's time to capture a frame (excluding the first frame)
        elif current_timestamp - prev_timestamp >= 1:
            # If there is a significant difference with the previous frame
            if prev_frame is not None and calculate_frame_difference(prev_frame, frame) > 10:
                # Construct the output file name
                output_filename = f"{output_folder}/frame_{int(current_timestamp):04d}.png"

                # Save the frame as an image
                cv2.imwrite(output_filename, frame)
                print(f"Saved {output_filename}")

                # Add the frame to the PowerPoint presentation
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                
                # Calculate the aspect ratio of the frame
                img_width, img_height = Image.open(output_filename).size
                aspect_ratio = img_width / img_height

                # Set the maximum width to almost fill the slide width
                max_picture_width = Inches(9.5)  # Adjust as needed based on your preference

                # Set the maximum height based on the aspect ratio
                max_picture_height = max_picture_width / aspect_ratio

                # Calculate the centered position
                slide_width, slide_height = presentation.slide_width, presentation.slide_height
                x_position = (slide_width - max_picture_width) / 2
                y_position = (slide_height - max_picture_height) / 2

                # Add the picture to the slide
                slide.shapes.add_picture(output_filename, x_position, y_position, width=max_picture_width, height=max_picture_height)

            # Update the previous frame and timestamp
            prev_frame = frame.copy()
            prev_timestamp = current_timestamp

    # Release the video capture object
    cap.release()

    # Delete the downloaded video file
    os.remove(video_path)
    print(f"Deleted the video file: {video_path}")

    # Save the PowerPoint presentation
    ppt_path = os.path.join(output_folder, f"{ppt_name}.pptx")
    presentation.save(ppt_path)
    print(f"Saved PowerPoint presentation: {ppt_path}")

if __name__ == "__main__":
    # Get user input for YouTube link
    youtube_video_url = input("Enter the YouTube video link: ")

    # Get user input for the folder name
    output_folder = input("Enter the folder name for screenshots: ")

    # Get user input for the PowerPoint presentation name
    ppt_name = input("Enter the PowerPoint presentation name: ")

    # Download the YouTube video with higher resolution
    video_path = download_youtube_video(youtube_video_url, output_folder, resolution="720p")

    # Call the capture_frames_with_transition function
    capture_frames_with_transition(video_path, output_folder, ppt_name)
